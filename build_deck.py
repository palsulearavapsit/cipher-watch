"""Build the RTMTS Scanner idea-submission deck from the OSF template.
Keeps the template frame (logos, slide numbers, titles) and fills 6 content
slides with polished content + screenshots + a flow diagram. Removes slide 7.
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy, os

TEMPLATE = r"C:\Users\HARSH\Downloads\OSF HACKONE PPT TEMPLATE 2k26.pptx"
OUT = r"C:\Users\HARSH\OneDrive\Desktop\osf hackathon\RTMTS_Scanner_Idea_Submission.pptx"
DL = r"C:\Users\HARSH\Downloads"
SHOT_DASH = os.path.join(DL, "WhatsApp Image 2026-06-23 at 8.23.19 PM.jpeg")
SHOT_LOGS = os.path.join(DL, "WhatsApp Image 2026-06-23 at 8.23.19 PM (1).jpeg")
SHOT_CHAIN = os.path.join(DL, "WhatsApp Image 2026-06-23 at 8.23.19 PM (2).jpeg")

# Palette (white-background deck)
NAVY   = "0E2233"   # titles
INK    = "1F2937"   # body
MUTED  = "6B7280"   # captions
CYAN   = "0E7490"   # brand accent (dark cyan = readable on white)
CYAN_L = "0891B2"
RED    = "DC2626"
AMBER  = "D97706"
GREEN  = "16A34A"
CHIP_BG= "E6F6F9"   # light cyan chip fill
CHIP_TX= "0E5A6B"
CARD_BG= "F4F7FA"   # subtle card fill

BODY_FONT = "Calibri"
TITLE_FONT = "Calibri"

prs = Presentation(TEMPLATE)
SW, SH = prs.slide_width, prs.slide_height

def find(slide, *subs, exclude_title=False):
    """Find first text shape whose text contains any of subs."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text
            if any(s.lower() in t.lower() for s in subs):
                return sh
    return None

def title_shape(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and (sh == slide.shapes.title or 'title' in sh.name.lower()):
            return sh
    return None

def style_runs(p, runs):
    for rt, o in runs:
        r = p.add_run(); r.text = rt
        f = r.font
        f.name = o.get('font', BODY_FONT)
        f.size = Pt(o.get('size', 12))
        f.bold = o.get('bold', False)
        f.italic = o.get('italic', False)
        col = o.get('color', INK)
        f.color.rgb = RGBColor.from_string(col)

def no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ('a:buChar', 'a:buAutoNum', 'a:buNone'):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn('a:buNone'), {}))

def write(tf, lines, anchor=None):
    """lines: list of dicts: {runs:[(text,opts)], sa, sb, align, level}"""
    tf.word_wrap = True
    tf.clear()
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(ln.get('sa', 6))
        p.space_before = Pt(ln.get('sb', 0))
        if 'line' in ln:
            p.line_spacing = ln['line']
        style_runs(p, ln['runs'])
        no_bullet(p)

def set_title(slide, text, size=32):
    ts = title_shape(slide)
    if ts is None:
        return
    write(ts.text_frame, [{'runs': [(text, {'size': size, 'bold': True, 'color': NAVY, 'font': TITLE_FONT})]}])
    # normalize title position (template has negative/odd offsets) & keep clear of top-right logos
    ts.left, ts.top, ts.width = Inches(0.55), Inches(0.95), Inches(9.0)

def add_box(slide, x, y, w, h, lines, fill=None, line=None, anchor=MSO_ANCHOR.TOP, pad=0.08):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = Inches(pad); tf.margin_right = Inches(pad)
    tf.margin_top = Inches(pad); tf.margin_bottom = Inches(pad)
    write(tf, lines, anchor=anchor)
    if fill:
        tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor.from_string(fill)
        if line:
            tb.line.color.rgb = RGBColor.from_string(line); tb.line.width = Pt(0.75)
        else:
            tb.line.fill.background()
    return tb

def add_card(slide, x, y, w, h, fill=CARD_BG, line=None, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp

def chip(slide, x, y, w, text, fill=CHIP_BG, tx=CHIP_TX, h=0.34, size=10.5):
    sp = add_card(slide, x, y, w, h, fill=fill, radius=0.5)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    write(tf, [{'runs': [(text, {'size': size, 'bold': True, 'color': tx})], 'align': PP_ALIGN.CENTER}],
          anchor=MSO_ANCHOR.MIDDLE)
    return sp

def add_pic(slide, path, x, y, w, h=None):
    return slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h) if h else None)

def arrow(slide, x, y, w=0.28, h=0.0, color=CYAN):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.18))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(color)
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

B = lambda t, o=None, **kw: (t, o if o is not None else kw)  # run helper (dict or kwargs)

# ---------------- SLIDE 1 — IDEA SUBMISSION ----------------
s = prs.slides[0]
heading = find(s, "IDEA SUBMISSION")
if heading:
    write(heading.text_frame, [{'runs': [("IDEA SUBMISSION", {'size': 30, 'bold': True, 'color': NAVY})]}])
body = find(s, "Theme", "Problem Statement", "Team")
if body:
    body.left, body.top, body.width, body.height = Inches(0.7), Inches(2.5), Inches(7.0), Inches(2.7)
    fields = [
        ("Theme", "[ Fill from portal ]"),
        ("Problem Statement ID", "[ Fill from portal ]"),
        ("Problem Statement Category", "[ Fill from portal ]"),
        ("Team ID", "[ Fill from portal ]"),
        ("Team Name (Registered on portal)", "[ Fill from portal ]"),
    ]
    lines = []
    for k, v in fields:
        lines.append({'runs': [B(k + " — ", size=14, bold=True, color=NAVY), B(v, size=14, color=MUTED)], 'sa': 8})
    write(body.text_frame, lines)
# product hero block (lower area)
add_box(s, 0.7, 5.25, 11.9, 1.6, [
    {'runs': [B("RTMTS Scanner", {'size': 40, 'bold': True, 'color': CYAN})], 'sa': 2},
    {'runs': [B("Real-Time Threat Monitoring System  ·  powered by the SentinelAI engine", {'size': 15, 'bold': True, 'color': INK})], 'sa': 4},
    {'runs': [B("Detects the threat. Scores the risk. Explains itself — in real time.", {'size': 14, 'italic': True, 'color': MUTED})]},
])

# ---------------- SLIDE 2 — IDEA TITLE / Solution Overview ----------------
s = prs.slides[1]
set_title(s, "IDEA TITLE", size=30)
body = find(s, "Solution Overview", "Description of your idea")
if body:
    body.left, body.top, body.width, body.height = Inches(0.4), Inches(2.18), Inches(7.55), Inches(4.0)
    write(body.text_frame, [
        {'runs': [B("RTMTS Scanner — a real-time, source-agnostic threat-detection platform that explains every alert.",
                    {'size': 15, 'bold': True, 'color': NAVY})], 'sa': 11, 'line': 1.08},
        {'runs': [B("WHAT IT IS", {'size': 12, 'bold': True, 'color': CYAN})], 'sa': 3},
        {'runs': [B("One engine watches every activity stream — payments, logins, database queries, blockchain wallets. It learns each entity's normal behavior and flags deviations in real time with a 0–100 risk score and a plain-English reason. A working prototype — 63 backend tests passing.",
                    {'size': 11.5, 'color': INK})], 'sa': 11, 'line': 1.1},
        {'runs': [B("HOW IT SOLVES THE PROBLEM", {'size': 12, 'bold': True, 'color': CYAN})], 'sa': 3},
        {'runs': [B("Legacy tools fire static-threshold alerts and dump raw rows with no context, so analysts drown in noise and catch attacks too late. We use per-entity behavioral baselines (catching novel attacks) and pre-explain every alert — triage drops from minutes to seconds.",
                    {'size': 11.5, 'color': INK})], 'sa': 11, 'line': 1.1},
        {'runs': [B("WHAT MAKES IT NOVEL", {'size': 12, 'bold': True, 'color': CYAN})], 'sa': 3},
        {'runs': [B("•  Source-agnostic — a new data source = one small adapter file. A platform, not a dashboard.", {'size': 11.5, 'color': INK})], 'sa': 5, 'line': 1.1},
        {'runs': [B("•  Explainable by default — Gemini turns an anomaly into one human sentence (cached, offline-safe).", {'size': 11.5, 'color': INK})], 'sa': 5, 'line': 1.1},
        {'runs': [B("•  Two-layer detection — statistical spike + ML drift fused into one intuitive score.", {'size': 11.5, 'color': INK})], 'sa': 5, 'line': 1.1},
    ])
# stat chip row (anchors the bottom-left)
stat_chips = [("63 tests passing", 2.35), ("Offline-safe demo", 2.45), ("Source-agnostic core", 2.55)]
cx = 0.45
for label, w in stat_chips:
    chip(s, cx, 6.35, w, label, h=0.40, size=11)
    cx += w + 0.18
# dashboard screenshot on the right (larger, vertically centered)
add_card(s, 8.05, 2.5, 5.12, 3.05, fill="FFFFFF", line="D7DEE6")
add_pic(s, SHOT_DASH, 8.18, 2.63, 4.86)
add_box(s, 8.05, 5.62, 5.12, 1.3, [
    {'runs': [B("LIVE SOC DASHBOARD", {'size': 11, 'bold': True, 'color': CYAN})], 'sa': 3, 'align': PP_ALIGN.CENTER},
    {'runs': [B("Press “inject attack” → within ~1–2s the entity is flagged, scored, and explained — every run, offline-safe.",
                {'size': 11, 'italic': True, 'color': MUTED})], 'align': PP_ALIGN.CENTER, 'line': 1.08},
])

# ---------------- SLIDE 3 — SOLUTION APPROACH ----------------
s = prs.slides[2]
set_title(s, "SOLUTION APPROACH", size=30)
body = find(s, "Technologies", "Methodology")
if body:
    # repurpose as the tech-stack header line
    body.left, body.top, body.width, body.height = Inches(0.55), Inches(1.78), Inches(12.2), Inches(0.4)
    write(body.text_frame, [{'runs': [B("TECHNOLOGY STACK", {'size': 12, 'bold': True, 'color': CYAN})]}])
# tech chips
techs = [
    ("React · Tailwind · Recharts", 2.45),
    ("FastAPI · WebSockets · Python", 2.65),
    ("scikit-learn IsolationForest", 2.55),
    ("Google Gemini", 1.5),
    ("Etherscan / Alchemy", 1.95),
]
x = 0.55
for label, w in techs:
    chip(s, x, 2.16, w, label)
    x += w + 0.16
# pipeline flow
steps = [("1  Ingest", "normalize"), ("2  Baseline", "per-entity"), ("3  Detect", "stat + ML"),
         ("4  Explain", "Gemini"), ("5  Stream", "WebSocket"), ("6  Act", "triage")]
fy = 2.78; bw = 1.78; bh = 0.72; gap = 0.30; fx = (13.33 - (bw*6 + gap*5)) / 2
add_box(s, 0.55, 2.62, 6.0, 0.3, [{'runs': [B("DETECTION PIPELINE", {'size': 12, 'bold': True, 'color': CYAN})]}])
fy = 2.96
for i, (t, sub) in enumerate(steps):
    cx = fx + i*(bw+gap)
    card = add_card(s, cx, fy, bw, bh, fill=CHIP_BG if i % 2 == 0 else "FFFFFF", line=CYAN, radius=0.12)
    write(card.text_frame, [
        {'runs': [B(t, {'size': 12, 'bold': True, 'color': CYAN})], 'align': PP_ALIGN.CENTER, 'sa': 0},
        {'runs': [B(sub, {'size': 9.5, 'color': MUTED})], 'align': PP_ALIGN.CENTER},
    ], anchor=MSO_ANCHOR.MIDDLE)
    if i < 5:
        arrow(s, cx + bw + 0.02, fy + bh/2 - 0.09)
# screenshots row (working prototype)
add_box(s, 0.55, 3.86, 8.0, 0.3, [{'runs': [B("WORKING PROTOTYPE", {'size': 12, 'bold': True, 'color': CYAN})]}])
shots = [(SHOT_DASH, "SOC Dashboard"), (SHOT_LOGS, "Threat Logs — scored events"), (SHOT_CHAIN, "Blockchain Monitor")]
sw = 3.92; sh_h = 2.13; sgap = 0.36; sx = (13.33 - (sw*3 + sgap*2)) / 2; sy = 4.18
for i, (path, cap) in enumerate(shots):
    cx = sx + i*(sw+sgap)
    add_card(s, cx, sy, sw, sh_h+0.04, fill="FFFFFF", line="D7DEE6")
    add_pic(s, path, cx+0.06, sy+0.06, sw-0.12)
    add_box(s, cx, sy+sh_h+0.08, sw, 0.3,
            [{'runs': [B(cap, {'size': 10.5, 'bold': True, 'color': INK})], 'align': PP_ALIGN.CENTER}])

# ---------------- SLIDE 4 — FEASIBILITY AND VIABILITY ----------------
s = prs.slides[3]
set_title(s, "FEASIBILITY AND VIABILITY", size=30)
body = find(s, "feasibility of the idea", "Potential challenges")
if body:
    body.left, body.top, body.width, body.height = Inches(0.55), Inches(2.12), Inches(5.55), Inches(4.9)
    write(body.text_frame, [
        {'runs': [B("FEASIBLE TODAY", {'size': 13, 'bold': True, 'color': CYAN})], 'sa': 6},
        {'runs': [B("•  Already built & demoable — live dashboard, real-time engine, attack-on-cue, plain-English explanations, 63 passing tests.", {'size': 12, 'color': INK})], 'sa': 6, 'line': 1.05},
        {'runs': [B("•  Lean 2-service architecture (React ↔ FastAPI) — runs on a single laptop, no cloud dependency in the demo.", {'size': 12, 'color': INK})], 'sa': 6, 'line': 1.05},
        {'runs': [B("•  Built entirely on proven, free open-source components — no exotic or costly infrastructure.", {'size': 12, 'color': INK})], 'sa': 6, 'line': 1.05},
    ])
# right column: risk -> mitigation card
add_card(s, 6.45, 2.12, 6.35, 4.55, fill=CARD_BG, radius=0.05)
add_box(s, 6.65, 2.26, 6.0, 0.4, [{'runs': [B("RISK  →  MITIGATION", {'size': 13, 'bold': True, 'color': NAVY})]}])
risks = [
    ("Live demo on bad wifi", "Gemini cached + template fallback; binds localhost — fully offline-safe."),
    ("ML cold-start (needs history)", "Warm-up seeding builds baselines; fast statistical layer carries the live spike."),
    ("False positives", "Per-entity baselines + 2-layer fusion; never learns from confirmed anomalies."),
    ("Scaling to real volume", "Stateless engine shards per entity; in-memory store is swappable."),
    ("Securing the product itself", "OWASP self-audit (SECURITY.md): no critical or high findings."),
]
ry = 2.74
for risk, mit in risks:
    add_box(s, 6.65, ry, 6.0, 0.74, [
        {'runs': [B(risk, {'size': 11.5, 'bold': True, 'color': RED})], 'sa': 1},
        {'runs': [B(mit, {'size': 10.5, 'color': INK})], 'line': 1.02},
    ])
    ry += 0.77

# ---------------- SLIDE 5 — IMPACT AND BENEFITS ----------------
s = prs.slides[4]
set_title(s, "IMPACT AND BENEFITS", size=30)
body = find(s, "impact on the target", "Benefits of the solution")
if body:
    body.left, body.top, body.width, body.height = Inches(0.55), Inches(2.05), Inches(12.2), Inches(0.55)
    write(body.text_frame, [{'runs': [
        B("Target audience:  ", {'size': 12.5, 'bold': True, 'color': NAVY}),
        B("banks & fintechs · SOC / security teams · crypto exchanges & wallets · any org with high-volume activity streams.",
          {'size': 12.5, 'color': INK})]}])
# stat callouts
stats = [("~1–2s", "from attack to flagged, scored & explained"),
         ("1 adapter", "to onboard a brand-new data source"),
         ("0", "critical / high security findings (OWASP audit)")]
cw = 3.92; cg = 0.36; cx0 = (13.33 - (cw*3 + cg*2)) / 2; cy = 2.7
for i, (big, lab) in enumerate(stats):
    cx = cx0 + i*(cw+cg)
    add_card(s, cx, cy, cw, 1.35, fill=CHIP_BG, radius=0.10)
    write(add_box(s, cx, cy, cw, 1.35, [], fill=None).text_frame, [
        {'runs': [B(big, {'size': 34, 'bold': True, 'color': CYAN})], 'align': PP_ALIGN.CENTER, 'sa': 2},
        {'runs': [B(lab, {'size': 11, 'color': INK})], 'align': PP_ALIGN.CENTER, 'line': 1.0},
    ], anchor=MSO_ANCHOR.MIDDLE) if False else None
    # write directly into the card
    c = add_box(s, cx, cy, cw, 1.35, [
        {'runs': [B(big, {'size': 34, 'bold': True, 'color': CYAN})], 'align': PP_ALIGN.CENTER, 'sa': 2},
        {'runs': [B(lab, {'size': 11, 'color': INK})], 'align': PP_ALIGN.CENTER, 'line': 1.0},
    ], anchor=MSO_ANCHOR.MIDDLE)
# benefits rows
add_box(s, 0.55, 4.3, 12.2, 0.34, [{'runs': [B("BENEFITS", {'size': 13, 'bold': True, 'color': CYAN})]}])
benefits = [
    ("Economic", "Cuts fraud losses, chargebacks & breach-remediation cost; fewer analyst-hours per alert; one platform replaces several point tools."),
    ("Social / trust", "Protects customers' money and data; faster response means less harm; supports safer digital finance."),
    ("Operational", "One engine for all sources = lower integration & maintenance cost; new source onboarded with a single adapter."),
    ("Efficiency", "Lean architecture with an in-memory hot path = low compute footprint vs. heavy traditional SIEM stacks."),
]
by = 4.68
for name, desc in benefits:
    add_box(s, 0.55, by, 12.2, 0.52, [{'runs': [
        B(name + " — ", {'size': 12, 'bold': True, 'color': NAVY}),
        B(desc, {'size': 11.5, 'color': INK})], 'line': 1.03}])
    by += 0.55

# ---------------- SLIDE 6 — RESEARCH AND REFERENCES ----------------
s = prs.slides[5]
set_title(s, "RESEARCH AND REFERENCES", size=30)
body = find(s, "Details / Links", "reference and research")
if body:
    body.left, body.top, body.width, body.height = Inches(0.55), Inches(2.1), Inches(12.2), Inches(4.6)
    refs = [
        ("Isolation Forest", "Liu, Ting & Zhou — “Isolation Forest,” IEEE ICDM 2008 (core anomaly-detection algorithm)."),
        ("scikit-learn", "IsolationForest & outlier-detection documentation — scikit-learn.org."),
        ("UEBA", "User & Entity Behavior Analytics — per-entity behavioral baselining (industry concept)."),
        ("OWASP Top 10 (2021)", "Framework used for our security self-audit — owasp.org/Top10."),
        ("Core stack docs", "FastAPI · React · Recharts · Tailwind CSS — official documentation."),
        ("Google Gemini API", "ai.google.dev — explanation layer."),
        ("Etherscan / Alchemy APIs", "blockchain wallet-activity module (optional)."),
        ("Validation roadmap", "Public fraud datasets (IEEE-CIS / credit-card fraud, Kaggle) for benchmarking."),
    ]
    lines = []
    for k, v in refs:
        lines.append({'runs': [B("•  " + k + " — ", {'size': 12, 'bold': True, 'color': CYAN}),
                               B(v, {'size': 12, 'color': INK})], 'sa': 7, 'line': 1.03})
    write(body.text_frame, lines)

# ---------------- remove slide 7 (Important Pointers) ----------------
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[6])

prs.save(OUT)
print("SAVED", OUT, "size_MB=", round(os.path.getsize(OUT)/1e6, 2), "slides=", len(prs.slides._sldIdLst))
